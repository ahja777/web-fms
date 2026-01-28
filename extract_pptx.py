# -*- coding: utf-8 -*-
from pptx import Presentation
import json
import os

pptx_path = r'C:\Claude_AI\backdata\A0203_G01_화면설계서(디지털물류네트워크)_Shipping_V1.0_20220718.pptx'
output_path = r'C:\Claude_AI\screen_design.json'

prs = Presentation(pptx_path)

reports = []
for idx, slide in enumerate(prs.slides, 1):
    slide_texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            slide_texts.append(shape.text.strip())
    if slide_texts:
        reports.append({'slide': idx, 'texts': slide_texts})

with open(output_path, 'w', encoding='utf-8') as f:
    json.dump(reports, f, ensure_ascii=False, indent=2)

print(f'Saved {len(reports)} slides to {output_path}')

# AWB 관련 페이지 출력
print("\n=== AWB/항공 관련 페이지 ===")
for slide in reports:
    all_text = " ".join(slide['texts']).upper()
    if "AWB" in all_text or "항공" in all_text or "AIR" in all_text or "WAYBILL" in all_text:
        print(f"\n슬라이드 {slide['slide']}:")
        for t in slide['texts']:
            print(f"  - {t[:150]}..." if len(t) > 150 else f"  - {t}")
